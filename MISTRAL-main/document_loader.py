"""Document loader and chunking module for office documents.

Supports PDF, DOCX, and TXT file formats with metadata extraction.
"""

import os
from pathlib import Path
from typing import Optional

from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document


def load_txt(file_path: str) -> str:
    """Load a plain text file."""
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def load_pdf(file_path: str) -> str:
    """Load a PDF file and extract text."""
    from pypdf import PdfReader

    reader = PdfReader(file_path)
    text_parts = []
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text_parts.append(page_text)
    return "\n\n".join(text_parts)


def load_docx(file_path: str) -> str:
    """Load a DOCX file and extract text."""
    from docx import Document as DocxDocument

    doc = DocxDocument(file_path)
    text_parts = []
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            text_parts.append(paragraph.text)
    return "\n\n".join(text_parts)


def load_xlsx(file_path: str) -> str:
    """Load an XLSX (Excel) file and extract text from all cells."""
    from openpyxl import load_workbook

    wb = load_workbook(file_path, data_only=True)
    text_parts = []
    
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        text_parts.append(f"Sheet: {sheet_name}")
        
        for row in sheet.iter_rows():
            row_text = []
            for cell in row:
                if cell.value is not None:
                    row_text.append(str(cell.value))
            if row_text:
                text_parts.append(" | ".join(row_text))
    
    return "\n".join(text_parts)


def load_pptx(file_path: str) -> str:
    """Load a PPTX (PowerPoint) file and extract text from all slides."""
    from pptx import Presentation

    prs = Presentation(file_path)
    text_parts = []
    
    for i, slide in enumerate(prs.slides, 1):
        text_parts.append(f"Slide {i}:")
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())
        
        text_parts.append("")  # Empty line between slides
    
    return "\n\n".join(text_parts)


LOADERS = {
    ".txt": load_txt,
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".xlsx": load_xlsx,
    ".pptx": load_pptx,
}


def load_document(file_path: str) -> Optional[str]:
    """Load a document based on its file extension.

    Returns None if the file type is unsupported.
    """
    ext = Path(file_path).suffix.lower()
    loader = LOADERS.get(ext)
    if loader is None:
        return None
    return loader(file_path)


def extract_metadata(file_path: str) -> dict:
    """Extract metadata from a file path."""
    path = Path(file_path)
    return {
        "source": str(path.resolve()),
        "filename": path.name,
        "file_type": path.suffix.lower(),
    }


def load_and_chunk_documents(
    documents_dir: str,
    chunk_size: int = 512,
    chunk_overlap: int = 50,
) -> list[Document]:
    """Load all supported documents from a directory and split into chunks.

    Args:
        documents_dir: Path to directory containing office documents.
        chunk_size: Maximum size of each text chunk.
        chunk_overlap: Overlap between consecutive chunks.

    Returns:
        List of LangChain Document objects with text and metadata.
    """
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", ". ", " ", ""],
    )

    all_chunks: list[Document] = []
    supported_extensions = set(LOADERS.keys())

    for root, _dirs, files in os.walk(documents_dir):
        for filename in files:
            file_path = os.path.join(root, filename)
            ext = Path(filename).suffix.lower()

            if ext not in supported_extensions:
                continue

            text = load_document(file_path)
            if not text or not text.strip():
                print(f"  Skipping empty file: {filename}")
                continue

            metadata = extract_metadata(file_path)
            chunks = splitter.create_documents(
                texts=[text],
                metadatas=[metadata],
            )
            all_chunks.extend(chunks)
            print(f"  Loaded {filename}: {len(chunks)} chunks")

    return all_chunks
